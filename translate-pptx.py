import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from googletrans import Translator
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import threading

translator = Translator(service_urls=['translate.googleapis.com'])

def translate_slides():
    # Get user input
    input_file = filedialog.askopenfilename(title="Select PowerPoint file")
    output_file = filedialog.asksaveasfilename(title="Save translated PowerPoint file", defaultextension=".pptx")
    destination_language = destination_language_entry.get()
    translation_mode = translation_mode_var.get()
    
    # Validate mode
    if not (translation_mode == 'merge' or translation_mode == 'overwrite'):
        status_label.config(text="Invalid mode. Please select 'merge' or 'overwrite'.")
        return
    
    # Create a separate thread for translation
    def translation_thread():
        try:
            prs = Presentation(input_file)
            num_slides = len(prs.slides)
            status_label.config(text=f"Translating {num_slides} slides...")
            i = 0
            for slide in prs.slides:
                i = i + 1
                shapes = []
                
                # only read in text-boxes
                for shape in slide.shapes:
                
                    if shape.has_text_frame and shape.text != '':
                        shapes.append(shape)
                
                status_label.config(text=f"Translating {i}/{len(prs.slides)} slides.")
                # Translate notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame
                    if translation_mode == 'merge':
                        notes.text = notes.text + '\n' + translator.translate(notes.text, dest=destination_language).text
                    elif translation_mode == 'overwrite':
                        notes.text = translator.translate(notes.text, dest=destination_language).text
                
                # Translate title
                if len(shapes) < 1:
                    continue
                
                title = shapes[0].text
                if translation_mode == 'merge':
                    shapes[0].text = title + '\n\n' + translator.translate(title, dest=destination_language).text
                elif translation_mode == 'overwrite':
                    shapes[0].text = translator.translate(title, dest=destination_language).text
                
                # Translate content
                for shape in shapes[1:]:
                    if shape.text != '':
                        content = shape.text
                        translation = translator.translate(content, dest=destination_language).text
                        
                        if translation_mode == 'merge':
                            p = shape.text_frame.add_paragraph()
                            p.alignment = PP_ALIGN.RIGHT
                            
                            run = p.add_run()
                            run.text = translation
                            font = run.font
                            run.font.language_id = 1025  # Arabic language ID
                            
                            font.size = Pt(10)
                        elif translation_mode == 'overwrite':
                            shape.text = translation
            
            prs.save(output_file)
            status_label.config(text=f"Translation completed. Saved as {output_file}")
        except Exception as e:
            status_label.config(text=f"An error occurred: {str(e)}")
    
    # Start the translation thread
    translation_thread = threading.Thread(target=translation_thread)
    translation_thread.start()

# Create the main window
window = tk.Tk()
window.title("PowerPoint Translator")

# Destination language selection
language_frame = tk.Frame(window)
language_frame.pack(pady=5)

language_label = tk.Label(language_frame, text="Destination Language:")
language_label.pack(side="left")

destination_language_entry = tk.Entry(language_frame)
destination_language_entry.pack(side="left")

# Translation mode selection
mode_frame = tk.Frame(window)
mode_frame.pack(pady=5)

translation_mode_label = tk.Label(mode_frame, text="Translation Mode:")
translation_mode_label.pack(side="left")

translation_mode_var = tk.StringVar(window, "merge")
merge_radio = tk.Radiobutton(mode_frame, text="Merge", variable=translation_mode_var, value="merge")
merge_radio.pack(side="left")

overwrite_radio = tk.Radiobutton(mode_frame, text="Overwrite", variable=translation_mode_var, value="overwrite")
overwrite_radio.pack(side="left")

# Input file selection
input_frame = tk.Frame(window)
input_frame.pack(pady=10)

input_label = tk.Label(input_frame, text="Select PowerPoint file:")
input_label.pack(side="left")

input_button = tk.Button(input_frame, text="Browse", command=translate_slides)
input_button.pack(side="left")


# Status label
status_label = tk.Label(window, text="")
status_label.pack(pady=10)

# Run the main window loop
window.mainloop()